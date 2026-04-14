#!/usr/bin/env python3
"""
AdStop Investor Presentation Generator
Clean & Professional design — 15 slides for pre-seed pitch
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Design Tokens ──────────────────────────────────────────────────────

# Colors
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT    = RGBColor(0xF8, 0xFA, 0xFC)
SURFACE     = RGBColor(0xF1, 0xF5, 0xF9)
BLUE        = RGBColor(0x25, 0x63, 0xEB)
BLUE_LIGHT  = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_DARK   = RGBColor(0x1E, 0x40, 0xAF)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
GREEN_LIGHT = RGBColor(0xD1, 0xFA, 0xE5)
TEXT_PRIMARY = RGBColor(0x0F, 0x17, 0x2A)
TEXT_SEC     = RGBColor(0x64, 0x74, 0x8B)
TEXT_LIGHT   = RGBColor(0x94, 0xA3, 0xB8)
BORDER       = RGBColor(0xE2, 0xE8, 0xF0)
RED          = RGBColor(0xEF, 0x44, 0x44)
RED_LIGHT    = RGBColor(0xFE, 0xE2, 0xE2)
AMBER        = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_LIGHT  = RGBColor(0xFE, 0xF3, 0xC7)
PURPLE       = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LIGHT = RGBColor(0xED, 0xE9, 0xFE)
CYAN         = RGBColor(0x06, 0xB6, 0xD4)
CYAN_LIGHT   = RGBColor(0xCF, 0xFA, 0xFE)
INDIGO       = RGBColor(0x63, 0x66, 0xF1)

FONT = "Calibri"

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Margins
MARGIN_L = Inches(0.8)
MARGIN_R = Inches(0.8)
MARGIN_T = Inches(0.6)
CONTENT_W = Inches(11.733)  # 13.333 - 0.8 - 0.8

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


# ── Helper Functions ───────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill=None, border_color=None, border_width=Pt(1)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill=None, border_color=None):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    # Reduce corner radius
    shape.adjustments[0] = 0.04
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=14, color=TEXT_PRIMARY,
             bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT, line_spacing=1.2):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, default_size=14,
                       default_color=TEXT_PRIMARY, alignment=PP_ALIGN.LEFT):
    """Add text box with multiple styled paragraphs.
    lines = [(text, size, color, bold, spacing_after), ...]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False
        spacing = line_data[4] if len(line_data) > 4 else 6

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT
        p.alignment = alignment
        p.space_after = Pt(spacing)
        p.space_before = Pt(0)
    return txBox


def add_accent_bar(slide, top=Inches(0), height=None):
    """Add a thin blue accent bar at the top."""
    if height is None:
        height = Inches(0.06)
    add_rect(slide, Inches(0), top, SLIDE_W, height, fill=BLUE)


def add_slide_number(slide, num, total=16):
    """Add slide number at bottom right."""
    add_text(slide, Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.3),
             f"{num}/{total}", font_size=10, color=TEXT_LIGHT, alignment=PP_ALIGN.RIGHT)


def add_logo_text(slide, left=MARGIN_L, top=Inches(0.25)):
    """Add AdStop text logo."""
    txBox = slide.shapes.add_textbox(left, top, Inches(2), Inches(0.45))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]

    run1 = p.add_run()
    run1.text = "Ad"
    run1.font.size = Pt(20)
    run1.font.color.rgb = TEXT_PRIMARY
    run1.font.bold = True
    run1.font.name = FONT

    run2 = p.add_run()
    run2.text = "Stop"
    run2.font.size = Pt(20)
    run2.font.color.rgb = BLUE
    run2.font.bold = True
    run2.font.name = FONT
    return txBox


def add_section_header(slide, title, subtitle=None, slide_num=1):
    """Standard section header with accent bar and title."""
    add_accent_bar(slide)
    add_logo_text(slide)
    add_text(slide, MARGIN_L, Inches(1.0), CONTENT_W, Inches(0.6),
             title, font_size=32, color=TEXT_PRIMARY, bold=True)
    if subtitle:
        add_text(slide, MARGIN_L, Inches(1.55), Inches(9.0), Inches(0.5),
                 subtitle, font_size=16, color=TEXT_SEC)
    add_slide_number(slide, slide_num)


def add_card(slide, left, top, width, height, title, body, icon="",
             title_size=16, body_size=12, fill=WHITE, border=BORDER,
             title_color=TEXT_PRIMARY, icon_color=BLUE):
    """Add a styled card with title and body text."""
    card = add_rounded_rect(slide, left, top, width, height, fill=fill, border_color=border)

    # Icon/emoji
    if icon:
        add_text(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.5), Inches(0.4),
                 icon, font_size=20, color=icon_color)
        t_left = left + Inches(0.2)
        t_top = top + Inches(0.55)
    else:
        t_left = left + Inches(0.2)
        t_top = top + Inches(0.15)

    # Title
    add_text(slide, t_left, t_top, width - Inches(0.4), Inches(0.35),
             title, font_size=title_size, color=title_color, bold=True)

    # Body
    add_text(slide, t_left, t_top + Inches(0.35), width - Inches(0.4), height - Inches(1.0),
             body, font_size=body_size, color=TEXT_SEC, line_spacing=1.4)

    return card


def add_stat_box(slide, left, top, width, height, number, label, sublabel="",
                 num_color=BLUE, bg_color=WHITE, border=BORDER):
    """Add a stat/metric box."""
    box = add_rounded_rect(slide, left, top, width, height, fill=bg_color, border_color=border)
    add_text(slide, left, top + Inches(0.15), width, Inches(0.5),
             number, font_size=36, color=num_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.65), width, Inches(0.3),
             label, font_size=13, color=TEXT_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    if sublabel:
        add_text(slide, left, top + Inches(0.95), width, Inches(0.3),
                 sublabel, font_size=10, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)
    return box


# ── SLIDE 1: Title ─────────────────────────────────────────────────────

def create_slide_1():
    slide = prs.slides.add_slide(blank_layout)

    # Full-width blue accent at top
    add_accent_bar(slide, height=Inches(0.08))

    # Left side blue decorative vertical bar
    add_rect(slide, Inches(0), Inches(0.08), Inches(0.08), Inches(7.42), fill=BLUE)

    # Main title area
    # "AdStop" in large text
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(8), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "Ad"
    run1.font.size = Pt(72)
    run1.font.color.rgb = TEXT_PRIMARY
    run1.font.bold = True
    run1.font.name = FONT
    run2 = p.add_run()
    run2.text = "Stop"
    run2.font.size = Pt(72)
    run2.font.color.rgb = BLUE
    run2.font.bold = True
    run2.font.name = FONT

    # Tagline
    add_text(slide, Inches(1.2), Inches(3.0), Inches(8), Inches(0.6),
             "Your Home Digital Shield", font_size=28, color=TEXT_SEC)

    # Description
    add_text(slide, Inches(1.2), Inches(3.7), Inches(8), Inches(0.8),
             "The plug-and-play device that blocks ads, trackers, and unwanted content\nacross your entire home network. No technical skills required.",
             font_size=16, color=TEXT_SEC, line_spacing=1.5)

    # Bottom stats bar
    bar_top = Inches(5.5)
    stats = [
        ("10,000+", "DNS queries blocked daily"),
        ("< 2 min", "Setup time"),
        ("25+", "Router brands supported"),
        ("2W", "Power consumption"),
    ]
    stat_w = Inches(2.6)
    for i, (num, label) in enumerate(stats):
        x = Inches(1.2) + i * (stat_w + Inches(0.3))
        add_rounded_rect(slide, x, bar_top, stat_w, Inches(1.2), fill=BG_LIGHT, border_color=BORDER)
        add_text(slide, x, bar_top + Inches(0.15), stat_w, Inches(0.45),
                 num, font_size=24, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(slide, x, bar_top + Inches(0.6), stat_w, Inches(0.4),
                 label, font_size=11, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)

    # Decorative circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(1.5), Inches(2.5), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE_LIGHT
    shape.line.fill.background()

    # Shield icon in the circle
    add_text(slide, Inches(10.8), Inches(2.0), Inches(2.0), Inches(1.5),
             "\U0001F6E1", font_size=60, color=BLUE, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 1)


# ── SLIDE 2: The Problem ──────────────────────────────────────────────

def create_slide_2():
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "The Problem", "Most households have no control over ads, trackers, and unwanted content.", 2)

    # Main problem statement
    add_rounded_rect(slide, MARGIN_L, Inches(2.2), CONTENT_W, Inches(1.0), fill=RED_LIGHT, border_color=None)
    add_text(slide, Inches(1.1), Inches(2.35), Inches(10.5), Inches(0.7),
             "Every connected device in your home is exposed to ads, trackers, and potentially harmful content \u2014 "
             "and most people have no way to stop it without technical expertise.",
             font_size=15, color=RGBColor(0x7F, 0x1D, 0x1D), line_spacing=1.4)

    # 4 problem cards
    cards = [
        ("\U0001F4FA", "Smart TV", "No ad-blocker available.\nAds play before every video,\nwith no way to skip or block."),
        ("\U0001F4F1", "Smartphone", "Ad-blockers only work in\none app or browser. Most\nin-app ads can't be blocked."),
        ("\U0001F3AE", "Console", "PlayStation, Xbox, Nintendo:\nimpossible to install any\nextension or ad-blocker."),
        ("\U0001F4BB", "PC & Tablet", "Requires setup on every\ndevice and every browser.\nFamily members forget."),
    ]

    card_w = Inches(2.7)
    card_h = Inches(2.8)
    gap = Inches(0.28)
    start_x = MARGIN_L
    card_y = Inches(3.6)

    for i, (icon, title, body) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rounded_rect(slide, x, card_y, card_w, card_h, fill=WHITE, border_color=BORDER)

        # Icon circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), card_y + Inches(0.2),
                                         Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BG_LIGHT
        circle.line.fill.background()
        add_text(slide, x + Inches(0.9), card_y + Inches(0.25), Inches(0.8), Inches(0.7),
                 icon, font_size=28, alignment=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.2), card_y + Inches(1.15), card_w - Inches(0.4), Inches(0.35),
                 title, font_size=16, color=TEXT_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.2), card_y + Inches(1.55), card_w - Inches(0.4), Inches(1.0),
                 body, font_size=12, color=TEXT_SEC, alignment=PP_ALIGN.CENTER, line_spacing=1.4)


# ── SLIDE 3: How It Works ─────────────────────────────────────────────

def create_slide_3():
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "How It Works", "Immediate setup, no technical skills required.", 3)

    steps = [
        ("01", "Connect", "Plug AdStop into power near your router. A setup network appears. Connect from your phone and enter your WiFi credentials via a simple captive portal.", BLUE, BLUE_LIGHT),
        ("02", "Auto-Configure", "AdStop connects to your network and configures itself as the DNS server. It auto-detects your router brand and applies optimal settings.", GREEN, GREEN_LIGHT),
        ("03", "Protected", "From that moment, every device on your network is automatically protected. No apps to install, no settings to change, nothing to maintain.", PURPLE, PURPLE_LIGHT),
    ]

    step_w = Inches(3.6)
    step_h = Inches(3.8)
    gap = Inches(0.4)
    start_x = MARGIN_L
    step_y = Inches(2.4)

    for i, (num, title, body, color, light_color) in enumerate(steps):
        x = start_x + i * (step_w + gap)

        # Card
        add_rounded_rect(slide, x, step_y, step_w, step_h, fill=WHITE, border_color=BORDER)

        # Step number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.35), step_y + Inches(0.3),
                                         Inches(0.9), Inches(0.9))
        circle.fill.solid()
        circle.fill.fore_color.rgb = light_color
        circle.line.fill.background()
        add_text(slide, x + Inches(1.35), step_y + Inches(0.38), Inches(0.9), Inches(0.7),
                 num, font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        # Title
        add_text(slide, x + Inches(0.3), step_y + Inches(1.4), step_w - Inches(0.6), Inches(0.35),
                 title, font_size=20, color=TEXT_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

        # Connector arrow (between steps)
        if i < 2:
            arrow_x = x + step_w + Inches(0.05)
            add_text(slide, arrow_x, step_y + Inches(1.5), Inches(0.3), Inches(0.4),
                     "\u2192", font_size=24, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

        # Body
        add_text(slide, x + Inches(0.3), step_y + Inches(1.9), step_w - Inches(0.6), Inches(1.6),
                 body, font_size=12, color=TEXT_SEC, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

    # Bottom note
    add_rounded_rect(slide, MARGIN_L, Inches(6.5), CONTENT_W, Inches(0.6), fill=GREEN_LIGHT)
    add_text(slide, Inches(1.1), Inches(6.55), Inches(10.5), Inches(0.4),
             "\u2713  Average setup time: under 2 minutes. Supports 25+ router brands including major Italian ISPs.",
             font_size=13, color=RGBColor(0x06, 0x5F, 0x46), bold=False)


# ── SLIDE 4: What's Inside ────────────────────────────────────────────

def create_slide_4():
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "What's Inside", "Compact, efficient, and built on proven open-source technology.", 4)

    # Left: Hardware description
    left_w = Inches(5.5)
    left_x = MARGIN_L

    add_rounded_rect(slide, left_x, Inches(2.3), left_w, Inches(4.5), fill=BG_LIGHT, border_color=BORDER)

    add_text(slide, left_x + Inches(0.3), Inches(2.45), left_w - Inches(0.6), Inches(0.35),
             "On-Board Hardware", font_size=20, color=TEXT_PRIMARY, bold=True)

    add_text(slide, left_x + Inches(0.3), Inches(2.85), left_w - Inches(0.6), Inches(1.2),
             "AdStop is built around a Raspberry Pi Zero 2W, a compact and energy-efficient "
             "single-board computer. It features a 2.13\" e-ink display that shows blocking "
             "statistics with near-zero power consumption.\n\n"
             "The software stack leverages proven open-source DNS filtering technology, "
             "enhanced with custom Rust networking components, a Python captive portal, "
             "and a proprietary cloud management system.",
             font_size=13, color=TEXT_SEC, line_spacing=1.5)

    # Components list
    components = [
        ("\u2022  Raspberry Pi Zero 2W (ARM Cortex-A53)", TEXT_PRIMARY),
        ("\u2022  Waveshare 2.13\" e-ink display (250x122)", TEXT_PRIMARY),
        ("\u2022  SanDisk 32GB microSD (Class A1)", TEXT_PRIMARY),
        ("\u2022  5V 2.5A power supply", TEXT_PRIMARY),
        ("\u2022  Custom 3D-printed PLA case", TEXT_PRIMARY),
    ]
    comp_y = Inches(4.6)
    for i, (text, color) in enumerate(components):
        add_text(slide, left_x + Inches(0.3), comp_y + Inches(i * 0.28), left_w - Inches(0.6), Inches(0.25),
                 text, font_size=12, color=color)

    # Right: Stats
    right_x = Inches(6.8)
    right_w = Inches(5.5)

    stats = [
        ("~35\u20ac", "Hardware Cost", "Per unit at volume\n(~55\u20ac prototype)", BLUE, BLUE_LIGHT),
        ("2W", "Power Draw", "Minimal energy impact\non your electricity bill", GREEN, GREEN_LIGHT),
        ("24/7", "Always On", "Silent operation,\nno maintenance needed", PURPLE, PURPLE_LIGHT),
    ]

    stat_h = Inches(1.3)
    for i, (num, label, sub, color, bg) in enumerate(stats):
        y = Inches(2.3) + i * (stat_h + Inches(0.15))
        add_rounded_rect(slide, right_x, y, right_w, stat_h, fill=WHITE, border_color=BORDER)

        # Colored accent left edge
        add_rect(slide, right_x + Inches(0.01), y + Inches(0.15), Inches(0.05), stat_h - Inches(0.3), fill=color)

        add_text(slide, right_x + Inches(0.35), y + Inches(0.15), Inches(1.5), Inches(0.5),
                 num, font_size=30, color=color, bold=True)
        add_text(slide, right_x + Inches(0.35), y + Inches(0.6), Inches(1.8), Inches(0.3),
                 label, font_size=14, color=TEXT_PRIMARY, bold=True)
        add_text(slide, right_x + Inches(2.3), y + Inches(0.25), Inches(2.8), Inches(0.8),
                 sub, font_size=11, color=TEXT_SEC, line_spacing=1.4)


# ── SLIDE 5: Why Different ────────────────────────────────────────────

def create_slide_5():
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "Why It's Different", "AdStop works at the network level, before ads even reach your devices.", 5)

    # Central diagram concept
    add_rounded_rect(slide, Inches(3.5), Inches(2.2), Inches(6.3), Inches(1.0), fill=BLUE_LIGHT, border_color=BLUE)
    add_text(slide, Inches(3.5), Inches(2.3), Inches(6.3), Inches(0.7),
             "\U0001F6E1  AdStop blocks unwanted DNS requests upstream \u2014 before they reach any device on your network.",
             font_size=14, color=BLUE_DARK, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    # Protected devices
    devices = [
        ("\U0001F4FA", "Smart TV", "No extensions possible.\nAdStop protects natively."),
        ("\U0001F3AE", "Game Console", "PS5, Xbox, Switch \u2014\nall automatically protected."),
        ("\U0001F4F1", "Smartphones", "Full protection without\ninstalling any app."),
        ("\U0001F4BB", "PCs & Tablets", "No per-device setup.\nOne device protects all."),
    ]

    card_w = Inches(2.7)
    card_h = Inches(2.5)
    gap = Inches(0.28)
    start_x = MARGIN_L + Inches(0.3)
    card_y = Inches(3.6)

    for i, (icon, title, body) in enumerate(devices):
        x = start_x + i * (card_w + gap)
        add_rounded_rect(slide, x, card_y, card_w, card_h, fill=WHITE, border_color=BORDER)
        add_text(slide, x, card_y + Inches(0.2), card_w, Inches(0.5),
                 icon, font_size=32, alignment=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.15), card_y + Inches(0.75), card_w - Inches(0.3), Inches(0.3),
                 title, font_size=15, color=TEXT_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.15), card_y + Inches(1.15), card_w - Inches(0.3), Inches(1.0),
                 body, font_size=11, color=TEXT_SEC, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    # vs traditional
    add_rounded_rect(slide, MARGIN_L, Inches(6.4), CONTENT_W, Inches(0.7), fill=BG_LIGHT)
    add_text(slide, Inches(1.1), Inches(6.45), Inches(10.5), Inches(0.5),
             "Traditional ad-blockers: per-device, per-browser, manual setup.    AdStop: one device, whole network, zero configuration.",
             font_size=13, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)


# ── SLIDE 6: What It Can Block ────────────────────────────────────────

def create_slide_6():
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "What It Can Block",
                       "A complete shield for the whole family, fully configurable from the app.", 6)

    categories = [
        ("\U0001F6AB", "Advertisements", "Banners, pre-roll videos, popups, and sponsored content across every platform and device.", RED, RED_LIGHT),
        ("\U0001F50D", "Trackers", "Tracking scripts that collect data about your online behavior, browsing habits, and personal info.", AMBER, AMBER_LIGHT),
        ("\U0001F3B0", "Gambling", "Betting sites and online casinos. Ideal for protecting younger family members.", PURPLE, PURPLE_LIGHT),
        ("\U0001F6E1", "Adult Content", "Configurable safe-browsing filter for the entire family. Enable or disable via the app.", BLUE, BLUE_LIGHT),
        ("\U0001F4F1", "Social Media", "Block access to social platforms to reduce distractions during study or work hours.", CYAN, CYAN_LIGHT),
    ]

    # Grid: 3 top, 2 bottom centered
    card_w = Inches(3.6)
    card_h = Inches(1.8)
    gap = Inches(0.3)

    for i, (icon, title, body, color, bg) in enumerate(categories):
        if i < 3:
            x = MARGIN_L + i * (card_w + gap)
            y = Inches(2.4)
        else:
            x = MARGIN_L + Inches(2.0) + (i - 3) * (card_w + gap)
            y = Inches(4.5)

        add_rounded_rect(slide, x, y, card_w, card_h, fill=WHITE, border_color=BORDER)

        # Color accent bar at top of card
        add_rect(slide, x + Inches(0.3), y + Inches(0.01), card_w - Inches(0.6), Inches(0.04), fill=color)

        add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.4),
                 icon, font_size=22)
        add_text(slide, x + Inches(0.7), y + Inches(0.2), card_w - Inches(1.0), Inches(0.3),
                 title, font_size=15, color=TEXT_PRIMARY, bold=True)
        add_text(slide, x + Inches(0.2), y + Inches(0.65), card_w - Inches(0.4), Inches(1.0),
                 body, font_size=11, color=TEXT_SEC, line_spacing=1.4)

    # Bottom note
    add_text(slide, MARGIN_L, Inches(6.6), CONTENT_W, Inches(0.4),
             "All categories can be enabled or disabled independently from the mobile app.",
             font_size=12, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


# ── SLIDE 7: The Control App ──────────────────────────────────────────

def create_slide_7():
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "The Control App",
                       "Complete management in the palm of your hand.", 7)

    # Status badge
    add_rounded_rect(slide, Inches(9.5), Inches(1.05), Inches(3.0), Inches(0.35),
                     fill=AMBER_LIGHT, border_color=None)
    add_text(slide, Inches(9.5), Inches(1.08), Inches(3.0), Inches(0.3),
             "\U0001F6A7  App in active development", font_size=11, color=AMBER, alignment=PP_ALIGN.CENTER)

    # Features list
    features = [
        ("\U0001F4CA", "Real-Time Statistics", "View how many ads and trackers have been blocked today, this week, and this month. Detailed charts and trends.", BLUE, BLUE_LIGHT),
        ("\u2699\uFE0F", "Filter Management", "Turn blocking categories on or off with a simple toggle. Customize protection levels per time of day.", GREEN, GREEN_LIGHT),
        ("\U0001F4F1", "Remote Control", "Manage AdStop even when you're away from home. Secure cloud connection with 2FA authentication.", PURPLE, PURPLE_LIGHT),
        ("\U0001F514", "Smart Alerts", "Get notified about security events, unusual traffic patterns, and device status changes.", RED, RED_LIGHT),
    ]

    card_w = Inches(5.5)
    card_h = Inches(1.15)
    gap = Inches(0.15)

    for i, (icon, title, body, color, bg) in enumerate(features):
        x = MARGIN_L
        y = Inches(2.4) + i * (card_h + gap)

        add_rounded_rect(slide, x, y, card_w, card_h, fill=WHITE, border_color=BORDER)
        add_rect(slide, x + Inches(0.01), y + Inches(0.15), Inches(0.05), card_h - Inches(0.3), fill=color)

        add_text(slide, x + Inches(0.25), y + Inches(0.1), Inches(0.4), Inches(0.35),
                 icon, font_size=18)
        add_text(slide, x + Inches(0.65), y + Inches(0.1), Inches(4.5), Inches(0.3),
                 title, font_size=14, color=TEXT_PRIMARY, bold=True)
        add_text(slide, x + Inches(0.65), y + Inches(0.45), Inches(4.5), Inches(0.6),
                 body, font_size=11, color=TEXT_SEC, line_spacing=1.3)

    # Right side: Platform info
    right_x = Inches(7.0)
    right_w = Inches(5.5)

    add_rounded_rect(slide, right_x, Inches(2.4), right_w, Inches(4.7), fill=BG_LIGHT, border_color=BORDER)

    add_text(slide, right_x + Inches(0.3), Inches(2.6), right_w - Inches(0.6), Inches(0.35),
             "Cross-Platform", font_size=18, color=TEXT_PRIMARY, bold=True)

    platforms = [
        ("\U0001F4F1  iOS", "Built with Capacitor \u2014 4 screens: Home, Activity, Alerts, Profile"),
        ("\U0001F4F1  Android", "Planned \u2014 same codebase via Capacitor"),
        ("\U0001F4BB  Desktop", "Electron app for macOS, Windows, Linux"),
        ("\U0001F310  Web", "Full Next.js dashboard accessible from any browser"),
    ]

    for i, (platform, desc) in enumerate(platforms):
        py = Inches(3.15) + i * Inches(0.85)
        add_rounded_rect(slide, right_x + Inches(0.2), py, right_w - Inches(0.4), Inches(0.7),
                         fill=WHITE, border_color=BORDER)
        add_text(slide, right_x + Inches(0.4), py + Inches(0.08), Inches(2.5), Inches(0.25),
                 platform, font_size=13, color=TEXT_PRIMARY, bold=True)
        add_text(slide, right_x + Inches(0.4), py + Inches(0.35), right_w - Inches(0.8), Inches(0.3),
                 desc, font_size=10, color=TEXT_SEC)


# ── SLIDE 8: Security ─────────────────────────────────────────────────

def create_slide_8():
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "Security & Updates",
                       "Security is built into every layer of AdStop, not bolted on as an afterthought.", 8)

    features = [
        ("\U0001F504", "Automatic OTA Updates",
         "Firmware updates are delivered over-the-air with automatic rollback if anything goes wrong. "
         "Users are always on the latest version without lifting a finger.",
         GREEN, GREEN_LIGHT),
        ("\U0001F511", "Signed Device License",
         "Each AdStop device has a digitally signed RSA-4096 license, ensuring authenticity. "
         "Factory provisioning prevents unauthorized cloning.",
         BLUE, BLUE_LIGHT),
        ("\U0001F512", "HMAC-Signed Communication",
         "All device-to-cloud communication is authenticated with HMAC signatures. "
         "Commands use double-HMAC for critical operations like remote lockdown.",
         PURPLE, PURPLE_LIGHT),
        ("\U0001F6E1", "TLS Encryption",
         "Cloud communication is secured via TLS. Full end-to-end encrypted channels "
         "are being deployed across all infrastructure components.",
         AMBER, AMBER_LIGHT),
    ]

    card_w = Inches(5.5)
    card_h = Inches(1.2)
    gap = Inches(0.12)

    for i, (icon, title, body, color, bg) in enumerate(features):
        col = i % 2
        row = i // 2
        x = MARGIN_L + col * (card_w + Inches(0.3))
        y = Inches(2.3) + row * (card_h + gap)

        add_rounded_rect(slide, x, y, card_w, card_h, fill=WHITE, border_color=BORDER)
        add_rect(slide, x + Inches(0.01), y + Inches(0.15), Inches(0.05), card_h - Inches(0.3), fill=color)

        add_text(slide, x + Inches(0.25), y + Inches(0.1), Inches(0.4), Inches(0.35),
                 icon, font_size=18)
        add_text(slide, x + Inches(0.65), y + Inches(0.1), card_w - Inches(0.9), Inches(0.3),
                 title, font_size=14, color=TEXT_PRIMARY, bold=True)
        add_text(slide, x + Inches(0.25), y + Inches(0.45), card_w - Inches(0.5), Inches(0.65),
                 body, font_size=10.5, color=TEXT_SEC, line_spacing=1.4)

    # Additional security features
    add_rounded_rect(slide, MARGIN_L, Inches(5.0), CONTENT_W, Inches(1.8), fill=BG_LIGHT, border_color=BORDER)
    add_text(slide, Inches(1.1), Inches(5.15), Inches(4.0), Inches(0.3),
             "Additional Security Measures", font_size=15, color=TEXT_PRIMARY, bold=True)

    extras = [
        "\u2022  Two-Factor Authentication (TOTP) for user accounts",
        "\u2022  Redis-backed distributed rate limiting on all endpoints",
        "\u2022  Complete security audit log (immutable event trail)",
        "\u2022  GDPR-compliant data export and deletion",
        "\u2022  Intrusion detection and real-time security monitoring",
    ]
    for i, text in enumerate(extras):
        add_text(slide, Inches(1.1), Inches(5.5) + Inches(i * 0.24), Inches(10.5), Inches(0.22),
                 text, font_size=11, color=TEXT_SEC)


# ── SLIDE 9: Project Status ───────────────────────────────────────────

def create_slide_9():
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "Project Status",
                       "AdStop is a working project with real code. Here\u2019s where we are today.", 9)

    components = [
        ("Cloud Backend", "85%", 0.85,
         "Complete Next.js web app: 2FA auth, MQTT device management, OTA updates, Stripe payments, "
         "admin panel, security audit logging, rate limiting, GDPR. Remaining: production deployment and monitoring.",
         BLUE, BLUE_LIGHT),
        ("Firmware", "55%", 0.55,
         "Built on Pi-hole DNS filtering + custom Rust networking components. Python captive portal, OTA with "
         "auto-rollback, auto-configuration for 25+ router brands. Blocker: rewriting first-boot system.",
         GREEN, GREEN_LIGHT),
        ("Mobile App", "50%", 0.50,
         "4 cross-platform screens via Capacitor (home, activity, alerts, profile). 2FA, onboarding, dark theme, "
         "Electron desktop version. Missing: Android build, push notifications, app store publication.",
         PURPLE, PURPLE_LIGHT),
        ("Hardware", "30%", 0.30,
         "Bill of materials finalized, 3D case model, wiring diagram. Missing: thermal testing, "
         "CE/FCC certifications, retail packaging for the v1 launch.",
         AMBER, AMBER_LIGHT),
    ]

    card_w = Inches(5.5)
    card_h = Inches(1.15)
    gap = Inches(0.12)

    for i, (name, pct, pct_val, desc, color, bg) in enumerate(components):
        col = i % 2
        row = i // 2
        x = MARGIN_L + col * (card_w + Inches(0.4))
        y = Inches(2.3) + row * (card_h + gap + Inches(0.05))

        add_rounded_rect(slide, x, y, card_w, card_h, fill=WHITE, border_color=BORDER)

        # Progress bar background
        bar_x = x + Inches(0.2)
        bar_y = y + Inches(0.12)
        bar_w = card_w - Inches(0.4)
        bar_h = Inches(0.16)
        add_rounded_rect(slide, bar_x, bar_y, bar_w, bar_h, fill=BG_LIGHT)
        # Progress bar fill
        if pct_val > 0:
            fill_w = int(bar_w * pct_val)
            add_rounded_rect(slide, bar_x, bar_y, fill_w, bar_h, fill=color)

        # Name and percentage
        add_text(slide, x + Inches(0.2), y + Inches(0.32), Inches(3.5), Inches(0.25),
                 name, font_size=14, color=TEXT_PRIMARY, bold=True)
        add_text(slide, x + Inches(4.0), y + Inches(0.32), Inches(1.2), Inches(0.25),
                 pct, font_size=14, color=color, bold=True, alignment=PP_ALIGN.RIGHT)

        # Description
        add_text(slide, x + Inches(0.2), y + Inches(0.6), card_w - Inches(0.4), Inches(0.5),
                 desc, font_size=9, color=TEXT_SEC, line_spacing=1.3)

    # Priority note
    add_rounded_rect(slide, MARGIN_L, Inches(5.0), CONTENT_W, Inches(0.7), fill=BLUE_LIGHT)
    add_text(slide, Inches(1.1), Inches(5.08), Inches(10.5), Inches(0.5),
             "\U0001F3AF  Immediate priority: complete the end-to-end flow \u2014 power on \u2192 setup \u2192 network protection live.",
             font_size=14, color=BLUE_DARK, bold=False)

    # Tech stack summary
    add_rounded_rect(slide, MARGIN_L, Inches(5.9), CONTENT_W, Inches(1.2), fill=BG_LIGHT, border_color=BORDER)
    add_text(slide, Inches(1.1), Inches(6.0), Inches(10.5), Inches(0.25),
             "Tech Stack", font_size=14, color=TEXT_PRIMARY, bold=True)

    stack_items = [
        "Cloud: Next.js 16, MySQL (Prisma), Redis, MQTT (Mosquitto), Stripe, Docker",
        "Firmware: Rust (DNS), Python (portal), Bash (provisioning), Pi-hole (filtering)",
        "App: Capacitor (iOS/Android), Electron (desktop), Tailwind CSS",
        "Security: Argon2id, HMAC-SHA256, RSA-4096, TOTP 2FA, Zod validation",
    ]
    for i, item in enumerate(stack_items):
        add_text(slide, Inches(1.1), Inches(6.3) + Inches(i * 0.2), Inches(10.5), Inches(0.2),
                 item, font_size=10, color=TEXT_SEC)


# ── SLIDE 10: Market Opportunity ───────────────────────────────────────

def create_slide_10():
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "Market Opportunity",
                       "A growing market with no dominant hardware solution for consumers.", 10)

    # Market stats
    stats = [
        ("42%", "of internet users\nuse ad-blockers", "Source: Statista 2024"),
        ("$100B+", "lost to ad fraud\nannually", "Growing concern for users"),
        ("0", "plug-and-play\nhardware solutions", "in the EU consumer market"),
    ]

    stat_w = Inches(3.6)
    stat_h = Inches(1.8)
    for i, (num, label, sub) in enumerate(stats):
        x = MARGIN_L + i * (stat_w + Inches(0.3))
        y = Inches(2.3)
        add_rounded_rect(slide, x, y, stat_w, stat_h, fill=WHITE, border_color=BORDER)
        add_text(slide, x, y + Inches(0.2), stat_w, Inches(0.5),
                 num, font_size=36, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(slide, x, y + Inches(0.75), stat_w, Inches(0.5),
                 label, font_size=13, color=TEXT_PRIMARY, alignment=PP_ALIGN.CENTER, line_spacing=1.3)
        add_text(slide, x, y + Inches(1.35), stat_w, Inches(0.3),
                 sub, font_size=10, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

    # Target segments
    add_text(slide, MARGIN_L, Inches(4.4), CONTENT_W, Inches(0.35),
             "Target Customers", font_size=18, color=TEXT_PRIMARY, bold=True)

    segments = [
        ("\U0001F468\u200D\U0001F469\u200D\U0001F467\u200D\U0001F466", "Families with Children",
         "Parents who want to protect kids from ads, gambling sites, and adult content without managing every device."),
        ("\U0001F512", "Privacy-Conscious Users",
         "People who care about tracking and data collection but don't have technical skills for DIY solutions."),
        ("\U0001F3E0", "Smart Home Owners",
         "Households with 10+ connected devices (TVs, consoles, cameras) that can't run traditional ad-blockers."),
        ("\U0001F468\u200D\U0001F4BB", "Remote Workers",
         "Professionals who want a distraction-free home network with social media and ad blocking during work hours."),
    ]

    seg_w = Inches(2.7)
    seg_h = Inches(2.2)
    gap = Inches(0.22)
    for i, (icon, title, body) in enumerate(segments):
        x = MARGIN_L + i * (seg_w + gap)
        y = Inches(4.9)
        add_rounded_rect(slide, x, y, seg_w, seg_h, fill=BG_LIGHT, border_color=BORDER)
        add_text(slide, x, y + Inches(0.15), seg_w, Inches(0.4),
                 icon, font_size=24, alignment=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.15), y + Inches(0.55), seg_w - Inches(0.3), Inches(0.3),
                 title, font_size=12, color=TEXT_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.15), y + Inches(0.9), seg_w - Inches(0.3), Inches(1.1),
                 body, font_size=10, color=TEXT_SEC, alignment=PP_ALIGN.CENTER, line_spacing=1.4)


# ── SLIDE 11: Business Model ──────────────────────────────────────────

def create_slide_11():
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "Business Model",
                       "Hardware sales with potential for recurring cloud subscription revenue.", 11)

    # Revenue streams
    add_text(slide, MARGIN_L, Inches(2.2), Inches(5.0), Inches(0.3),
             "Revenue Streams", font_size=18, color=TEXT_PRIMARY, bold=True)

    streams = [
        ("\U0001F4E6", "Hardware Sales", "Direct-to-consumer device sales. Hardware cost ~35\u20ac at volume, "
         "with healthy margin potential at consumer pricing (TBD).", BLUE, BLUE_LIGHT),
        ("\U0001F504", "Cloud Subscription", "Optional premium tier for advanced features: detailed analytics, "
         "custom blocklists, remote management, priority support.", GREEN, GREEN_LIGHT),
        ("\U0001F3ED", "B2B / Bulk", "Office and co-working space packages. Multi-device management "
         "with centralized admin dashboard.", PURPLE, PURPLE_LIGHT),
    ]

    card_w = Inches(5.5)
    card_h = Inches(1.1)
    for i, (icon, title, body, color, bg) in enumerate(streams):
        y = Inches(2.65) + i * (card_h + Inches(0.12))
        add_rounded_rect(slide, MARGIN_L, y, card_w, card_h, fill=WHITE, border_color=BORDER)
        add_rect(slide, MARGIN_L + Inches(0.01), y + Inches(0.12), Inches(0.05), card_h - Inches(0.24), fill=color)
        add_text(slide, MARGIN_L + Inches(0.25), y + Inches(0.1), Inches(0.4), Inches(0.3), icon, font_size=18)
        add_text(slide, MARGIN_L + Inches(0.65), y + Inches(0.1), Inches(4.5), Inches(0.25),
                 title, font_size=14, color=TEXT_PRIMARY, bold=True)
        add_text(slide, MARGIN_L + Inches(0.65), y + Inches(0.4), Inches(4.5), Inches(0.6),
                 body, font_size=10.5, color=TEXT_SEC, line_spacing=1.4)

    # Distribution channels
    right_x = Inches(7.0)
    right_w = Inches(5.5)

    add_text(slide, right_x, Inches(2.2), right_w, Inches(0.3),
             "Distribution Channels", font_size=18, color=TEXT_PRIMARY, bold=True)

    channels = [
        ("Direct Online", "Company website with Stripe checkout. Lowest CAC, highest margin."),
        ("Amazon / Marketplaces", "Reach millions of customers. Higher fees but massive exposure."),
        ("Electronics Retail", "Partnership with retailers for physical shelf presence."),
        ("Tech Communities", "Pi-hole / privacy forums, tech blogs, YouTube reviewers."),
    ]

    for i, (title, desc) in enumerate(channels):
        y = Inches(2.65) + i * Inches(0.85)
        add_rounded_rect(slide, right_x, y, right_w, Inches(0.72), fill=BG_LIGHT, border_color=BORDER)
        add_text(slide, right_x + Inches(0.25), y + Inches(0.08), right_w - Inches(0.5), Inches(0.22),
                 title, font_size=13, color=TEXT_PRIMARY, bold=True)
        add_text(slide, right_x + Inches(0.25), y + Inches(0.34), right_w - Inches(0.5), Inches(0.3),
                 desc, font_size=10, color=TEXT_SEC)

    # Unit economics note
    add_rounded_rect(slide, MARGIN_L, Inches(6.3), CONTENT_W, Inches(0.7), fill=GREEN_LIGHT)
    add_text(slide, Inches(1.1), Inches(6.4), Inches(10.5), Inches(0.4),
             "\U0001F4B0  Unit economics: hardware BOM ~35\u20ac at volume | Pricing strategy under evaluation | "
             "Stripe billing infrastructure already integrated",
             font_size=12, color=RGBColor(0x06, 0x5F, 0x46))


# ── SLIDE 12: Competitive Landscape ───────────────────────────────────

def create_slide_12():
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "Competitive Landscape",
                       "Where AdStop fits in the ad-blocking ecosystem.", 12)

    # Comparison table using cards
    competitors = [
        ("Pi-hole", "Open-source DNS sinkhole", [
            ("Setup", "Complex \u2014 requires Raspberry Pi, terminal, SSH", RED),
            ("Users", "Technical hobbyists only", AMBER),
            ("App", "Web interface only, no mobile app", AMBER),
            ("Cloud", "None \u2014 local only", RED),
            ("Hardware", "DIY assembly required", RED),
        ]),
        ("NextDNS", "Cloud DNS filtering", [
            ("Setup", "Easy \u2014 change DNS settings", GREEN),
            ("Users", "General consumers", GREEN),
            ("App", "Mobile app available", GREEN),
            ("Cloud", "Cloud-only, no local control", AMBER),
            ("Hardware", "None \u2014 software only", RED),
        ]),
        ("AdGuard Home", "Self-hosted DNS filter", [
            ("Setup", "Moderate \u2014 needs server/NAS", AMBER),
            ("Users", "Tech-savvy users", AMBER),
            ("App", "Web interface, basic mobile", AMBER),
            ("Cloud", "Optional cloud sync", AMBER),
            ("Hardware", "DIY or NAS required", AMBER),
        ]),
        ("AdStop", "Plug-and-play hardware", [
            ("Setup", "2 minutes \u2014 plug in and go", GREEN),
            ("Users", "Everyone \u2014 zero tech skills", GREEN),
            ("App", "Full mobile + desktop app", GREEN),
            ("Cloud", "Integrated cloud management", GREEN),
            ("Hardware", "Dedicated device included", GREEN),
        ]),
    ]

    card_w = Inches(2.75)
    card_h = Inches(4.2)
    gap = Inches(0.2)

    for i, (name, desc, features) in enumerate(competitors):
        x = MARGIN_L + i * (card_w + gap)
        y = Inches(2.3)

        is_adstop = name == "AdStop"
        border = BLUE if is_adstop else BORDER
        bg = BLUE_LIGHT if is_adstop else WHITE

        add_rounded_rect(slide, x, y, card_w, card_h, fill=bg, border_color=border)

        if is_adstop:
            add_rect(slide, x, y, card_w, Inches(0.05), fill=BLUE)

        # Name
        name_color = BLUE if is_adstop else TEXT_PRIMARY
        add_text(slide, x + Inches(0.15), y + Inches(0.15), card_w - Inches(0.3), Inches(0.3),
                 name, font_size=18, color=name_color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.15), y + Inches(0.45), card_w - Inches(0.3), Inches(0.25),
                 desc, font_size=9, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)

        # Features
        for j, (label, value, color) in enumerate(features):
            fy = y + Inches(0.85) + j * Inches(0.62)
            add_text(slide, x + Inches(0.15), fy, card_w - Inches(0.3), Inches(0.18),
                     label, font_size=9, color=TEXT_LIGHT, bold=True)
            add_text(slide, x + Inches(0.15), fy + Inches(0.18), card_w - Inches(0.3), Inches(0.35),
                     value, font_size=10, color=color)

    # Positioning statement
    add_rounded_rect(slide, MARGIN_L, Inches(6.7), CONTENT_W, Inches(0.55), fill=BLUE_LIGHT)
    add_text(slide, Inches(1.1), Inches(6.75), Inches(10.5), Inches(0.4),
             "AdStop is the only product that combines dedicated hardware + cloud management + mobile app "
             "in a zero-configuration consumer package.",
             font_size=12, color=BLUE_DARK, alignment=PP_ALIGN.CENTER)


# ── SLIDE 13: Roadmap ─────────────────────────────────────────────────

def create_slide_13():
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "Roadmap", "From working prototype to market launch.", 13)

    phases = [
        ("Q2 2026", "MVP Complete", [
            "End-to-end flow working",
            "Cloud platform deployed",
            "First-boot system rewritten",
            "10 beta units assembled",
        ], BLUE, BLUE_LIGHT, True),
        ("Q3 2026", "Beta Program", [
            "50 users beta testing",
            "Mobile app on TestFlight",
            "Android build ready",
            "Iterate on feedback",
        ], GREEN, GREEN_LIGHT, False),
        ("Q4 2026", "Pre-Orders", [
            "Custom PCB design finalized",
            "CE marking process started",
            "Pre-order campaign launch",
            "First production batch (100 units)",
        ], PURPLE, PURPLE_LIGHT, False),
        ("2027", "Scale", [
            "App Store & Play Store launch",
            "Amazon marketplace listing",
            "B2B enterprise offering",
            "1,000+ units target",
        ], AMBER, AMBER_LIGHT, False),
    ]

    card_w = Inches(2.75)
    card_h = Inches(4.2)
    gap = Inches(0.2)

    for i, (period, title, items, color, bg, current) in enumerate(phases):
        x = MARGIN_L + i * (card_w + gap)
        y = Inches(2.3)

        border = color if current else BORDER
        add_rounded_rect(slide, x, y, card_w, card_h, fill=WHITE, border_color=border)

        # Top color bar
        add_rect(slide, x + Inches(0.15), y + Inches(0.01), card_w - Inches(0.3), Inches(0.05), fill=color)

        # Current badge
        if current:
            add_rounded_rect(slide, x + Inches(0.6), y + Inches(0.15), Inches(1.55), Inches(0.28),
                             fill=bg, border_color=None)
            add_text(slide, x + Inches(0.6), y + Inches(0.15), Inches(1.55), Inches(0.25),
                     "WE ARE HERE", font_size=9, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        # Period
        period_y = y + Inches(0.55) if current else y + Inches(0.25)
        add_text(slide, x + Inches(0.15), period_y, card_w - Inches(0.3), Inches(0.3),
                 period, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.15), period_y + Inches(0.35), card_w - Inches(0.3), Inches(0.3),
                 title, font_size=14, color=TEXT_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

        # Items
        items_y = period_y + Inches(0.8)
        for j, item in enumerate(items):
            add_text(slide, x + Inches(0.2), items_y + j * Inches(0.35), card_w - Inches(0.4), Inches(0.3),
                     f"\u2022  {item}", font_size=11, color=TEXT_SEC)

        # Arrow between phases
        if i < 3:
            arrow_x = x + card_w + Inches(0.02)
            add_text(slide, arrow_x, Inches(3.8), Inches(0.2), Inches(0.4),
                     "\u203A", font_size=20, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


# ── SLIDE 14: Future Upgrades & Vision ─────────────────────────────────

def create_slide_14_future():
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "Future Upgrades & Vision",
                       "AdStop v1 is just the beginning. Here's where we're heading after launch.", 14)

    upgrades = [
        ("\U0001F50C", "Custom PCB Design",
         "Replace the Raspberry Pi with a custom-designed board. Target: cut production cost from "
         "~35\u20ac to ~15\u201320\u20ac per unit, enable a slimmer form factor, and secure the supply chain.",
         BLUE, BLUE_LIGHT),
        ("\U0001F9E0", "AI Threat Detection",
         "Anomaly detection for suspicious traffic, phishing, and zero-day domains. Network-wide "
         "telemetry creates a data moat that strengthens with every new device sold.",
         PURPLE, PURPLE_LIGHT),
        ("\U0001F310", "Built-in VPN",
         "One-click network-wide VPN for the entire household. Opens a large adjacent market with "
         "recurring subscription revenue and stronger privacy positioning.",
         GREEN, GREEN_LIGHT),
        ("\U0001F3E0", "Smart Home Integration",
         "Native integration with Apple HomeKit, Matter, Google Home, and Alexa. Voice control, "
         "automations, and seamless ecosystem interoperability.",
         CYAN, CYAN_LIGHT),
        ("\U0001F465", "Parental Controls Pro",
         "Screen time scheduling per device, homework mode, bedtime mode, and location-based rules. "
         "A premium subscription tier targeted at families.",
         AMBER, AMBER_LIGHT),
        ("\U0001F3E2", "Enterprise Tier",
         "B2B offering for offices, schools, and co-working spaces. Multi-location management, "
         "SSO/SAML, compliance reporting \u2014 much higher ACV per customer.",
         RED, RED_LIGHT),
    ]

    card_w = Inches(3.75)
    card_h = Inches(2.05)
    gap_x = Inches(0.23)
    gap_y = Inches(0.2)

    for i, (icon, title, body, color, bg) in enumerate(upgrades):
        col = i % 3
        row = i // 3
        x = MARGIN_L + col * (card_w + gap_x)
        y = Inches(2.3) + row * (card_h + gap_y)

        add_rounded_rect(slide, x, y, card_w, card_h, fill=WHITE, border_color=BORDER)

        # Left color accent bar
        add_rect(slide, x + Inches(0.01), y + Inches(0.15), Inches(0.05), card_h - Inches(0.3), fill=color)

        # Icon circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.22), y + Inches(0.22),
                                         Inches(0.65), Inches(0.65))
        circle.fill.solid()
        circle.fill.fore_color.rgb = bg
        circle.line.fill.background()
        add_text(slide, x + Inches(0.22), y + Inches(0.27), Inches(0.65), Inches(0.55),
                 icon, font_size=20, alignment=PP_ALIGN.CENTER)

        # Title
        add_text(slide, x + Inches(0.95), y + Inches(0.32), card_w - Inches(1.1), Inches(0.35),
                 title, font_size=14, color=color, bold=True)

        # Body
        add_text(slide, x + Inches(0.22), y + Inches(1.0), card_w - Inches(0.4), Inches(1.0),
                 body, font_size=10, color=TEXT_SEC, line_spacing=1.4)

    # Vision statement at bottom
    add_rounded_rect(slide, MARGIN_L, Inches(6.85), CONTENT_W, Inches(0.45), fill=BLUE_LIGHT)
    add_text(slide, Inches(1.1), Inches(6.9), Inches(10.5), Inches(0.35),
             "\U0001F680  Our long-term vision: AdStop becomes the default privacy hub in every home \u2014 not just an ad blocker.",
             font_size=12, color=BLUE_DARK, alignment=PP_ALIGN.CENTER)


# ── SLIDE 15: Team ─────────────────────────────────────────────────────

def create_slide_14():
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "The Team", "A small, focused team building something real.", 15)

    members = [
        ("\U0001F468\u200D\U0001F4BB", "Founder & Lead Developer",
         "Full-stack architecture, cloud platform, firmware development. "
         "Designed and built the entire AdStop system from hardware to cloud.",
         ["System Architecture", "Next.js / React", "Rust", "IoT / MQTT", "Security"],
         BLUE, BLUE_LIGHT),
        ("\U0001F527", "Co-Founder / Hardware",
         "Hardware design, PCB prototyping, 3D case modeling, "
         "manufacturing planning and supply chain.",
         ["Electronics Design", "3D Modeling", "Manufacturing", "Supply Chain"],
         GREEN, GREEN_LIGHT),
        ("\U0001F3A8", "Design & UX",
         "User interface design, mobile app UX, branding, "
         "and marketing materials.",
         ["UI/UX Design", "Mobile Design", "Branding", "Marketing"],
         PURPLE, PURPLE_LIGHT),
    ]

    card_w = Inches(3.6)
    card_h = Inches(4.0)
    gap = Inches(0.3)
    start_x = MARGIN_L + Inches(0.5)

    for i, (icon, role, desc, skills, color, bg) in enumerate(members):
        x = start_x + i * (card_w + gap)
        y = Inches(2.3)

        add_rounded_rect(slide, x, y, card_w, card_h, fill=WHITE, border_color=BORDER)

        # Icon circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.3), y + Inches(0.2),
                                         Inches(1.0), Inches(1.0))
        circle.fill.solid()
        circle.fill.fore_color.rgb = bg
        circle.line.fill.background()
        add_text(slide, x + Inches(1.3), y + Inches(0.25), Inches(1.0), Inches(0.9),
                 icon, font_size=36, alignment=PP_ALIGN.CENTER)

        # Role
        add_text(slide, x + Inches(0.2), y + Inches(1.3), card_w - Inches(0.4), Inches(0.35),
                 role, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        add_text(slide, x + Inches(0.2), y + Inches(1.7), card_w - Inches(0.4), Inches(0.9),
                 desc, font_size=11, color=TEXT_SEC, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

        # Skills tags
        skills_y = y + Inches(2.7)
        tag_x = x + Inches(0.15)
        for j, skill in enumerate(skills):
            row = j // 2
            col = j % 2
            tw = (card_w - Inches(0.4)) / 2
            tx = x + Inches(0.2) + col * tw
            ty = skills_y + row * Inches(0.32)
            add_rounded_rect(slide, tx, ty, tw - Inches(0.05), Inches(0.28), fill=bg, border_color=None)
            add_text(slide, tx, ty + Inches(0.02), tw - Inches(0.05), Inches(0.22),
                     skill, font_size=9, color=color, alignment=PP_ALIGN.CENTER)

    # Bottom note
    add_text(slide, MARGIN_L, Inches(6.6), CONTENT_W, Inches(0.4),
             "Made with \u2764 in Italy",
             font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


# ── SLIDE 16: Contact / Next Steps ────────────────────────────────────

def create_slide_15():
    slide = prs.slides.add_slide(blank_layout)

    # Blue accent bar
    add_accent_bar(slide, height=Inches(0.08))
    add_rect(slide, Inches(0), Inches(0.08), Inches(0.08), Inches(7.42), fill=BLUE)

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(8), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "Let's Build the Future\nof Home Privacy"
    run1.font.size = Pt(40)
    run1.font.color.rgb = TEXT_PRIMARY
    run1.font.bold = True
    run1.font.name = FONT

    # Subtitle
    add_text(slide, Inches(1.2), Inches(2.9), Inches(8), Inches(0.6),
             "We're raising a pre-seed round to bring AdStop from working prototype to market.",
             font_size=18, color=TEXT_SEC)

    # Use of funds
    add_text(slide, Inches(1.2), Inches(3.8), Inches(5.0), Inches(0.35),
             "Use of Funds", font_size=18, color=TEXT_PRIMARY, bold=True)

    funds = [
        ("\u2022  Complete product development (MVP \u2192 v1.0)", "\u2022  First production batch (100 units)"),
        ("\u2022  CE/FCC certifications", "\u2022  App Store publication (iOS + Android)"),
        ("\u2022  Initial marketing & pre-orders", "\u2022  Operational runway (6 months)"),
    ]

    for i, (left, right) in enumerate(funds):
        fy = Inches(4.25) + i * Inches(0.32)
        add_text(slide, Inches(1.2), fy, Inches(5.0), Inches(0.28),
                 left, font_size=13, color=TEXT_SEC)
        add_text(slide, Inches(6.5), fy, Inches(5.0), Inches(0.28),
                 right, font_size=13, color=TEXT_SEC)

    # Contact card
    contact_y = Inches(5.4)
    add_rounded_rect(slide, Inches(1.2), contact_y, Inches(5.5), Inches(1.4),
                     fill=BG_LIGHT, border_color=BORDER)

    # AdStop logo text in contact
    txBox2 = slide.shapes.add_textbox(Inches(1.5), contact_y + Inches(0.15), Inches(3), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    r1 = p2.add_run()
    r1.text = "Ad"
    r1.font.size = Pt(22)
    r1.font.color.rgb = TEXT_PRIMARY
    r1.font.bold = True
    r1.font.name = FONT
    r2 = p2.add_run()
    r2.text = "Stop"
    r2.font.size = Pt(22)
    r2.font.color.rgb = BLUE
    r2.font.bold = True
    r2.font.name = FONT

    contacts = [
        "\U0001F4E7  support@adstop.io",
        "\U0001F310  adstop.io",
    ]
    for i, contact in enumerate(contacts):
        add_text(slide, Inches(1.5), contact_y + Inches(0.6) + i * Inches(0.3),
                 Inches(4.5), Inches(0.25),
                 contact, font_size=13, color=TEXT_SEC)

    # Decorative circle (matching slide 1)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.0), Inches(3.5), Inches(3.0), Inches(3.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE_LIGHT
    shape.line.fill.background()
    add_text(slide, Inches(10.3), Inches(4.1), Inches(2.5), Inches(1.8),
             "\U0001F6E1", font_size=72, color=BLUE, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 16)


# ── Generate All Slides ───────────────────────────────────────────────

print("Generating AdStop investor presentation...")

create_slide_1()
print("  [1/16] Title slide")
create_slide_2()
print("  [2/16] The Problem")
create_slide_3()
print("  [3/16] How It Works")
create_slide_4()
print("  [4/16] What's Inside")
create_slide_5()
print("  [5/16] Why It's Different")
create_slide_6()
print("  [6/16] What It Can Block")
create_slide_7()
print("  [7/16] The Control App")
create_slide_8()
print("  [8/16] Security & Updates")
create_slide_9()
print("  [9/16] Project Status")
create_slide_10()
print("  [10/16] Market Opportunity")
create_slide_11()
print("  [11/16] Business Model")
create_slide_12()
print("  [12/16] Competitive Landscape")
create_slide_13()
print("  [13/16] Roadmap")
create_slide_14_future()
print("  [14/16] Future Upgrades & Vision")
create_slide_14()
print("  [15/16] Team")
create_slide_15()
print("  [16/16] Contact / Next Steps")

output_path = "/Users/marwansalah/Desktop/AdstopDocs/AdStop.pptx"
prs.save(output_path)
print(f"\nDone! Saved to {output_path}")
print(f"Backup of original: /Users/marwansalah/Desktop/AdstopDocs/AdStop_backup.pptx")
